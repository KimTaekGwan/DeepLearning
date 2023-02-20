from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pdfminer.high_level import extract_text
import PyPDF2
import fitz

from collections import defaultdict
from kiwipiepy import Kiwi

import re
import os
import shutil
from bs4 import BeautifulSoup

class DBUpdater:
    def __init__(self) -> None:


class Util:
    def __init__(self) -> None:
        self.input_dir = 'data/input/'
        self.res_dir = 'data/image/'
        self.orignal_dir = 'data/original/'
        
    def input_files_update(self):
        for file_name in os.listdir(self.input_dir):
            self.dir_Update(file_name)
    
    def dir_Update(self, file_name):
        name, ext = os.path.splitext(file_name)
        input_path = os.path.join(self.input_dir, file_name)
        img_path = os.path.join(self.res_dir, name)
        data_path = os.path.join(self.orignal_dir, file_name)
        
        if os.path.isfile(data_path):
            with open(input_path, 'rb') as input_file, open(data_path, 'rb') as data_file:
                input_content = input_file.read()
                data_content = data_file.read()
                if input_content == data_content:
                    print(f"{file_name} already exists")
                else:
                    # file 이름 변경 후 다시 재귀하기
                    print(f"{file_name} 이름이 같은 파일이 존재 - 이름을 변경해서 올려주세요")
                    # os.rename()
        else:
            os.makedirs(img_path)
            shutil.move(input_path, self.orignal_dir)
    

class Text_Preprocessing:
    def __init__(self) -> None:
        self.punct = "/-'?!.,#$%\'()*+-/:;<=>@[\\]^_`{|}~" + '""“”’' + '∞θ÷α•à−β∅³π‘₹´°£€\×™√²—–&'
        self.mapping = {"‘": "'", "₹": "e", "´": "'", "°": "", 
                                "€": "e", "™": "tm", "√": " sqrt ", "×": "x", "²": "2", 
                                "—": "-", "–": "-", "’": "'", "_": "-", "`": "'", '“': '"', '”': '"', 
                                '“': '"', "£": "e", '∞': 'infinity', 'θ': 'theta', '÷': '/', 'α': 'alpha', 
                                '•': '.', 'à': 'a', '−': '-', 'β': 'beta', '∅': '', '³': '3', 'π': 'pi', } 
        self.specials = {'\u200b': ' ', '…': ' ... ', '\ufeff': '', 'करना': '', 'है': ''}

    def preprocessing(self, text):
        text = self.clean(text)
        text = self.clean_str(text)
        text = self.token(text)
        return text
    
    def clean(self, text):
        for p in self.mapping:
            text = text.replace(p, self.mapping[p])
        
        for p in self.punct:
            text = text.replace(p, f' {p} ')
        
        
        for s in self.specials:
            text = text.replace(s, self.specials[s])
        return text.strip()

    def clean_str(self, text):
        # pattern = '([a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+)' # E-mail제거
        # text = re.sub(pattern=pattern, repl='', string=text)
        pattern = '(http|ftp|https)://(?:[-\w.]|(?:%[\da-fA-F]{2}))+' # URL제거
        text = re.sub(pattern=pattern, repl='', string=text)
        pattern = '([ㄱ-ㅎㅏ-ㅣ]+)'  # 한글 자음, 모음 제거
        text = re.sub(pattern=pattern, repl='', string=text)
        pattern = '<[^>]*>'         # HTML 태그 제거
        text = re.sub(pattern=pattern, repl='', string=text)
        pattern = '[^\w\s\n]'         # 특수기호제거
        text = re.sub(pattern=pattern, repl='', string=text)
        text = re.sub('[-=+,#/\?:^$.@*\"※~&%ㆍ!』\\‘|\(\)\[\]\<\>`\'…》]',' ', string=text)
        text = re.sub('\n', ' ', string=text)
        return text 
    
    def token(self, text) -> list:
        ls = [t.strip() for t in text.split(' ') if t]
        res = ' '.join(ls)
        return res


class Extract_Info:
    def __init__(self) -> None:
        self.tp = Text_Preprocessing()
        self.util = Util()
        self.kiwi = Kiwi()
        self.reset()
    
    def reset(self):
        self.file_name = None
        self.name = None
        self.ext = None
        self.text_dict = {'total':[], 'page':defaultdict(list)}
        self.link_list = []
        self.img_dict = {}
        self.text_info = None
        self.page_num = None
        self.img_num = 0
        self.image_blob = []
        self.pre_text = None
        
        
class PDF_Info_Extract(Extract_Info):
    def __init__(self) -> None:
        super().__init__()
        
    def reset(self):
        super().reset()
        
    
    def extract(self, file_name):
        self.reset()
        self.file_name = file_name
        self.name, self.ext = os.path.splitext(file_name)
        
        self.text_extract(file_name)
        self.image_extract(file_name)
        
        self.last()
        self.preprocessing()
    
    def last(self):
        res = [t for t in self.text_dict['total'][0].split('\n') if t]
        self.text_info = '\n'.join(res)
    
    def preprocessing(self):
        self.pre_text = ' '.join([t for t in self.text_info.split('\n') if t])
        
        self.pre_text = self.kiwi.space(self.pre_text)
        # print(1)
        # print(self.pre_text)
        self.pre_text = BeautifulSoup(self.pre_text, 'html.parser').text 
        # print(2)
        # print(self.pre_text)
        self.pre_text = re.sub(r'[^ ㄱ-ㅣ가-힣]', '', self.pre_text) #특수기호 제거, 정규 표현식
        
        # print(3)
        clean_words = []
        for token, pos, _, _ in self.kiwi.analyze(self.pre_text)[0][0]:
            if pos.startswith('N'):
                clean_words.append(token)
        # print(clean_words)
        self.pre_text = ' '.join(clean_words)
    
    def text_extract(self, file_name):
        path = self.util.orignal_dir + file_name
        
        try:
            text = extract_text(path)
        except:
            with open(path, 'rb') as f:
                # Create a PDF reader object
                pdf_reader = PyPDF2.PdfReader(f)
                # Get the total number of pages
                num_pages = pdf_reader.pages
                # Loop through each page and extract the text
                text = ''
                for page in num_pages:
                    text += page.extract_text()
        
        pattern = 'http[s]?://(?:[a-zA-Z]|[0-9]|[$\-@\.&+:/?=]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        links = re.findall(pattern, text)
        if links:
            self.link_list = links
            text = re.sub(pattern=pattern, repl='', string=text)
        
        if text:
            self.text_dict['total'].append(text)
    
    def image_extract(self, file_name):
        path = self.util.orignal_dir + file_name
        open_file = fitz.open(path)
        
        for page_number in range(len(open_file)):
            page = open_file[page_number]
            list_image = page.get_images()
            
            if list_image:
                # print(f"{len(list_image)} images found on page {page_number}")
                for image in list_image:
                    # Get the pixmap of the image
                    pixmap = fitz.Pixmap(open_file, image[0])
                    if pixmap.height >= 200 and pixmap.width >= 200:
                        # self.img_num  - img_num
                        self.img_num += 1
                        num = str(self.img_num).zfill(3)
                        image_filename = f"data/image/{self.name}/image_{self.img_num}.jpeg"
                        pixmap.save(image_filename)
            # else:
            #     print("No images found on page", page_number)
        open_file.close()


class PPT_Info_Extract(Extract_Info):
    def __init__(self) -> None:
        super().__init__()
    
    def reset(self):
        super().reset()
        self.parsed = None
    
    def extract(self, file_name):
        # path = "deep_learning_intro.pptx"
        self.reset()
        self.file_name = file_name
        self.name, self.ext = os.path.splitext(file_name)
        path = self.util.orignal_dir + file_name
        self.parsed = Presentation(path)
        
        for idx, slide in enumerate(self.parsed.slides):
            self.page_num = idx+1
            for shape in slide.shapes:
                self.group_check(shape)
        
        self.last()
        self.preprocessing()
    
    def last(self):
        res = []
        for text_ls in self.text_dict['page'].values():
            ls = []
            for text in text_ls:
                # print(text)
                pattern = 'http[s]?://(?:[a-zA-Z]|[0-9]|[$\-@\.&+:/?=]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
                links = re.findall(pattern, text)
                if links:
                    for link in links:
                        self.link_list.append(link)
                    text = re.sub(pattern=pattern, repl='', string=text)
                
                if text:
                    ls.append(text)
                    self.text_dict['total'].append(text)
            t = self.tp.preprocessing(' '.join(ls))
            if t:
                res.append(t)
        
        self.text_info = '\n'.join(res)
        
    def preprocessing(self):
        self.pre_text = ' '.join([t for t in self.text_info.split('\n') if t])
        
        self.pre_text = self.kiwi.space(self.pre_text)
        # print(1)
        # print(self.pre_text)
        self.pre_text = BeautifulSoup(self.pre_text, 'html.parser').text 
        # print(2)
        # print(self.pre_text)
        self.pre_text = re.sub(r'[^ ㄱ-ㅣ가-힣]', '', self.pre_text) #특수기호 제거, 정규 표현식
        
        # print(3)
        clean_words = []
        for token, pos, _, _ in self.kiwi.analyze(self.pre_text)[0][0]:
            if pos.startswith('N'):
                clean_words.append(token)
        # print(clean_words)
        # print(clean_words)
        self.pre_text = ' '.join(clean_words)
    
    def group_check(self, shape):
        # 그룹 모형
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self.info_extract(s)
        # 그룹 모형 아님
        else:
            self.info_extract(shape)
    
    def info_extract(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            self.group_check(shape)
        else:
            # yes image
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_blob = shape.image.blob
                
                # 이미지 중복 체크
                if image_blob not in self.image_blob:
                    ext = shape.image.ext # - (확장명)
                    size = shape.image.size
                    if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                        self.img_num += 1
                        num = str(self.img_num).zfill(3)
                        save_path =f"data/image/{self.name}/image_{num}.{ext}"
                        with open(save_path, "wb") as file:
                            file.write(image_blob)
                        self.image_blob.append(image_blob)
                        self.img_dict[self.img_num] = save_path
            # no image
            else:
                # yes text
                if shape.has_text_frame:
                    if shape.text.strip() != "":
                        self.text_dict['page'][self.page_num].append(shape.text.strip())